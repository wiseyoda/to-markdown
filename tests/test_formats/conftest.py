"""Session-scoped fixture factories for format-specific tests.

Each fixture generates a real file using the appropriate library,
cached per test session for performance.
"""

from pathlib import Path

import pytest

# ---------------------------------------------------------------------------
# PDF fixtures (fpdf2)
# ---------------------------------------------------------------------------


@pytest.fixture(scope="session")
def pdf_simple(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """Single-page PDF with plain text."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(text="Hello World. This is a simple PDF document.")
    path = tmp_path_factory.mktemp("pdf") / "simple.pdf"
    pdf.output(str(path))
    return path


@pytest.fixture(scope="session")
def pdf_with_headings(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """Multi-page PDF with heading-like text and paragraphs."""
    from fpdf import FPDF

    pdf = FPDF()
    # Page 1
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.multi_cell(w=0, text="Chapter One")
    pdf.ln(5)
    pdf.set_font("Helvetica", size=12)
    pdf.multi_cell(
        w=0,
        text="This is the first chapter of the document. "
        "It contains introductory material about the topic.",
    )
    # Page 2
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.multi_cell(w=0, text="Chapter Two")
    pdf.ln(5)
    pdf.set_font("Helvetica", size=12)
    pdf.multi_cell(
        w=0,
        text="The second chapter continues with more detailed analysis. "
        "It builds on concepts from the first chapter.",
    )
    path = tmp_path_factory.mktemp("pdf") / "headings.pdf"
    pdf.output(str(path))
    return path


@pytest.fixture(scope="session")
def pdf_with_table(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """PDF containing a simple 3x3 table."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.multi_cell(w=0, text="Table Example")
    pdf.ln(5)

    headers = ["Name", "Age", "City"]
    rows = [
        ["Alice", "30", "New York"],
        ["Bob", "25", "London"],
        ["Carol", "35", "Tokyo"],
    ]
    col_width = 50
    row_height = 10

    pdf.set_font("Helvetica", "B", 12)
    for header in headers:
        pdf.cell(w=col_width, h=row_height, text=header, border=1)
    pdf.ln(row_height)

    pdf.set_font("Helvetica", size=12)
    for row in rows:
        for cell in row:
            pdf.cell(w=col_width, h=row_height, text=cell, border=1)
        pdf.ln(row_height)

    path = tmp_path_factory.mktemp("pdf") / "table.pdf"
    pdf.output(str(path))
    return path


# ---------------------------------------------------------------------------
# DOCX fixtures (python-docx)
# ---------------------------------------------------------------------------


@pytest.fixture(scope="session")
def docx_simple(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """DOCX with a single paragraph of text."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("This is a simple document with one paragraph of text.")
    path = tmp_path_factory.mktemp("docx") / "simple.docx"
    doc.save(str(path))
    return path


@pytest.fixture(scope="session")
def docx_with_formatting(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """DOCX with headings, bold, italic, and a bullet list."""
    from docx import Document

    doc = Document()
    doc.add_heading("Main Title", level=1)
    doc.add_heading("Section One", level=2)
    para = doc.add_paragraph()
    para.add_run("This paragraph has ").bold = False
    bold_run = para.add_run("bold text")
    bold_run.bold = True
    para.add_run(" and ")
    italic_run = para.add_run("italic text")
    italic_run.italic = True
    para.add_run(" in it.")

    doc.add_heading("Section Two", level=2)
    doc.add_paragraph("First item", style="List Bullet")
    doc.add_paragraph("Second item", style="List Bullet")
    doc.add_paragraph("Third item", style="List Bullet")
    path = tmp_path_factory.mktemp("docx") / "formatted.docx"
    doc.save(str(path))
    return path


@pytest.fixture(scope="session")
def docx_with_table(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """DOCX containing a 3-column table with headers and data."""
    from docx import Document

    doc = Document()
    doc.add_heading("Employee Directory", level=1)
    table = doc.add_table(rows=4, cols=3)
    table.style = "Table Grid"

    headers = ["Name", "Department", "Role"]
    for i, header in enumerate(headers):
        table.rows[0].cells[i].text = header

    data = [
        ["Alice", "Engineering", "Developer"],
        ["Bob", "Marketing", "Designer"],
        ["Carol", "Engineering", "Manager"],
    ]
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, value in enumerate(row_data):
            table.rows[row_idx].cells[col_idx].text = value

    path = tmp_path_factory.mktemp("docx") / "table.docx"
    doc.save(str(path))
    return path


# ---------------------------------------------------------------------------
# PPTX fixtures (python-pptx)
# ---------------------------------------------------------------------------


@pytest.fixture(scope="session")
def pptx_simple(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """PPTX with a single slide containing title and body text."""
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Welcome Slide"
    slide.placeholders[1].text = "This is the body text of the presentation."
    path = tmp_path_factory.mktemp("pptx") / "simple.pptx"
    prs.save(str(path))
    return path


@pytest.fixture(scope="session")
def pptx_multi_slide(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """PPTX with 3 slides; slide 2 has speaker notes."""
    from pptx import Presentation

    prs = Presentation()

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Overview"
    slide1.placeholders[1].text = "Q1 2026 Update"

    # Slide 2: Content with speaker notes
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Metrics"
    slide2.placeholders[1].text = "Revenue increased by 15% year over year."
    notes_slide = slide2.notes_slide
    notes_slide.notes_text_frame.text = "Remember to highlight the growth trend."

    # Slide 3: Content
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Next Steps"
    slide3.placeholders[1].text = "Focus on customer retention and expansion."

    path = tmp_path_factory.mktemp("pptx") / "multi_slide.pptx"
    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# XLSX fixtures (openpyxl)
# ---------------------------------------------------------------------------


@pytest.fixture(scope="session")
def xlsx_simple(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """XLSX with a single sheet: header row + 3 data rows."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Product", "Price", "Quantity"])
    ws.append(["Widget", 9.99, 100])
    ws.append(["Gadget", 24.99, 50])
    ws.append(["Doohickey", 4.99, 200])
    path = tmp_path_factory.mktemp("xlsx") / "simple.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture(scope="session")
def xlsx_multi_sheet(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """XLSX with two sheets: Sales and Expenses."""
    from openpyxl import Workbook

    wb = Workbook()

    ws_sales = wb.active
    ws_sales.title = "Sales"
    ws_sales.append(["Month", "Revenue", "Units"])
    ws_sales.append(["January", 10000, 150])
    ws_sales.append(["February", 12000, 180])

    ws_expenses = wb.create_sheet("Expenses")
    ws_expenses.append(["Category", "Amount"])
    ws_expenses.append(["Rent", 3000])
    ws_expenses.append(["Utilities", 500])
    ws_expenses.append(["Supplies", 750])

    path = tmp_path_factory.mktemp("xlsx") / "multi_sheet.xlsx"
    wb.save(str(path))
    return path


# ---------------------------------------------------------------------------
# HTML fixtures (static strings)
# ---------------------------------------------------------------------------


_HTML_SIMPLE = """\
<!DOCTYPE html>
<html>
<head><title>Simple Page</title></head>
<body>
<h1>Welcome</h1>
<p>This is a simple HTML page with a heading and paragraph.</p>
</body>
</html>
"""

_HTML_WITH_TABLE = """\
<!DOCTYPE html>
<html>
<head><title>Table Page</title></head>
<body>
<h1>Data Report</h1>
<h2>Summary</h2>
<p>Below is the quarterly data.</p>
<table>
  <tr><th>Quarter</th><th>Revenue</th><th>Growth</th></tr>
  <tr><td>Q1</td><td>$100K</td><td>5%</td></tr>
  <tr><td>Q2</td><td>$120K</td><td>20%</td></tr>
  <tr><td>Q3</td><td>$115K</td><td>-4%</td></tr>
</table>
<h2>Links</h2>
<ul>
  <li><a href="https://example.com">Example Site</a></li>
  <li><a href="https://example.org">Another Site</a></li>
</ul>
</body>
</html>
"""


@pytest.fixture(scope="session")
def html_simple(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """HTML page with title, heading, and paragraph."""
    path = tmp_path_factory.mktemp("html") / "simple.html"
    path.write_text(_HTML_SIMPLE, encoding="utf-8")
    return path


@pytest.fixture(scope="session")
def html_with_table(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """HTML page with headings, table, and links."""
    path = tmp_path_factory.mktemp("html") / "table.html"
    path.write_text(_HTML_WITH_TABLE, encoding="utf-8")
    return path


# ---------------------------------------------------------------------------
# Image fixtures (Pillow)
# ---------------------------------------------------------------------------


@pytest.fixture(scope="session")
def image_png_simple(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """200x100 white PNG with 'Test Image' text."""
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (200, 100), color="white")
    draw = ImageDraw.Draw(img)
    draw.text((10, 40), "Test Image", fill="black")
    path = tmp_path_factory.mktemp("images") / "text.png"
    img.save(str(path))
    return path


@pytest.fixture(scope="session")
def image_jpg_simple(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """200x100 colored JPEG."""
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (200, 100), color="steelblue")
    draw = ImageDraw.Draw(img)
    draw.rectangle([20, 20, 180, 80], fill="gold", outline="black")
    path = tmp_path_factory.mktemp("images") / "colored.jpg"
    img.save(str(path), "JPEG")
    return path
